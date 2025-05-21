from utils.file_utils import load_csv
from pptx import Presentation
from io import BytesIO
# from data_preparation import load_csv
from datetime import datetime
from datetime import date
from dateutil.relativedelta import relativedelta
from pptx.util import Inches, Pt
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
import pandas as pd
import matplotlib.pyplot as plt
import io

def table_summary_channel (df, selected_period):
    df['month'] = pd.to_datetime(df['month'])  # ubah ke datetime
    df = df.sort_values('month')  # pastikan urut
    df['month'] = df['month'].dt.strftime('%Y-%m')

    df['followers_growth'] = df['followers'].diff()
    df['total_post_growth'] = (df['total_post'].diff()/ df['total_post'].shift(1) * 100).round(0).fillna(0).astype(int)
    df['reach_growth'] = (df['reach'].diff()/ df['reach'].shift(1) * 100).round(0).fillna(0).astype(int)
    df['engagement_growth'] = (df['engagement'].diff() / df['engagement'].shift(1) * 100).round(0).fillna(0).astype(int)
    df['profile_visit_growth'] = (df['profile_visit'].diff() / df['profile_visit'].shift(1) * 100).round(0).fillna(0).astype(int)

    df['post_save_rate'] = df['saved_count'] / df['total_post']
    df['CTR'] = (df['profile_visit'] / df['profile_reach']) * 100
    # Ambil hanya untuk bulan terpilih
    df_filtered = df[df['month'] == selected_period]
    # Output final
    final_df = df_filtered[[
        "month",
        "followers", "followers_growth",
        "total_post", "total_post_growth",
        "reach", "reach_growth",
        "engagement", "engagement_growth",
        "profile_visit", "profile_visit_growth",
        "post_save_rate", "CTR"
    ]]
    return(final_df)
def table_followers_trend(df, selected_period):
    selected_year = selected_period[:4]  # Ambil "2025" dari "2025-03"

    # Filter dari bulan Januari sampai selected_period di tahun itu
    filtered = df[(df['month'] >= f"{selected_year}-01") & (df['month'] <= selected_period)]

    return filtered[['month', 'followers']].sort_values('month')

def table_engagement_trend(df, selected_period):
    selected_year = selected_period[:4]  # Ambil "2025" dari "2025-03"

    # Filter dari bulan Januari sampai selected_period di tahun itu
    filtered = df[(df['month'] >= f"{selected_year}-01") & (df['month'] <= selected_period)]

    return filtered[['month', 'engagement']].sort_values('month')

def table_top_post (df, selected_period):
    df['post_date'] = pd.to_datetime(df['post_date'])
    df['month'] = df['post_date'].dt.strftime('%Y-%m')
    df['engagement'] = df['like_count'] + df['comment_count'] + df['share_count'] + df['saved_count']
    filtered_df = df[df['month'] == selected_period]
    top_posts = filtered_df.sort_values(by='engagement', ascending=False)
    top_1_posts = top_posts.head(1)
    top_1_posts = top_1_posts[['media_url','engagement', 'media_type']]
    return top_1_posts
def generate_text_top_post(df):
    if df.empty:
        return "No top post found."

    row = df.iloc[0]
    return (
        f"Engagement: {row['engagement']}\n"
        f"Media type: {row['media_type']}\n"
        f"URL: {row['media_url']}"
    )
def table_bottom_post (df, selected_period):
    df['post_date'] = pd.to_datetime(df['post_date'])
    df['month'] = df['post_date'].dt.strftime('%Y-%m')
    df['engagement'] = df['like_count'] + df['comment_count'] + df['share_count'] + df['saved_count']
    filtered_df = df[df['month'] == selected_period]
    top_posts = filtered_df.sort_values(by='engagement', ascending=True)
    top_1_posts = top_posts.head(1)
    top_1_posts = top_1_posts[['media_url','engagement', 'media_type']]
    return(top_1_posts)
def table_popular_post(df, selected_period):
    df['post_date'] = pd.to_datetime(df['post_date'])
    df['month'] = df['post_date'].dt.strftime('%Y-%m')
    df['engagement'] = df['like_count'] + df['comment_count'] + df['share_count'] + df['saved_count']
    # Ekstrak tahun dari selected_period
    selected_year = selected_period[:4]
    # Filter dari awal tahun hingga selected_period
    filtered_df = df[(df['month'] >= f"{selected_year}-01") & (df['month'] <= selected_period)]
    top_posts = filtered_df.sort_values(by='engagement', ascending=False)
    top_1_posts = top_posts.head(1)
    top_1_posts = top_1_posts[['media_url','engagement', 'media_type']]
    return(top_1_posts)


def table_evaluation_media(df, selected_period):
    selected_period = datetime.strptime(selected_period, "%Y-%m")
    previous_period = selected_period - relativedelta(months=1)
    filtered_df = df[df['month'].isin([previous_period.strftime('%Y-%m'),
                                       selected_period.strftime('%Y-%m')])]

    previous_period_name = previous_period.strftime('%B %Y')
    selected_period_name = selected_period.strftime('%B %Y')

    before_data = filtered_df[filtered_df['month'] ==
                              previous_period.strftime('%Y-%m')].iloc[0,1:].to_dict()
    last_data = filtered_df[filtered_df['month'] ==
                            selected_period.strftime('%Y-%m')].iloc[0,1:].to_dict()

    # Create a new DataFrame with Metric, Before, Last columns
    evaluation_media = pd.DataFrame(columns=['Metric', previous_period_name, selected_period_name])
    for metric, before_value in before_data.items():
        last_value = last_data.get(metric)
        evaluation_media = pd.concat([evaluation_media, pd.DataFrame(
            {'Metric': [metric], previous_period_name: [before_value], selected_period_name: [last_value]})],
                                     ignore_index=True)

    return evaluation_media


def table_evaluation_socmed(df, selected_period):
    df = df.drop(columns=['like_count','comment_count','saved_count','share_count','growth','profile_visit','profile_reach'])
    selected_period = datetime.strptime(selected_period, "%Y-%m")
    previous_period = selected_period - relativedelta(months=1)
    filtered_df = df[df['month'].isin([previous_period.strftime('%Y-%m'),
                                       selected_period.strftime('%Y-%m')])]

    previous_period_name = previous_period.strftime('%B %Y')
    selected_period_name = selected_period.strftime('%B %Y')

    # Get unique channels from the filtered DataFrame
    channels = filtered_df['channel'].unique()

    # Create a list to store evaluation DataFrames for each channel
    all_channel_evaluations = []

    # Iterate through each channel
    for channel in channels:
        # Filter data for the current channel
        channel_df = filtered_df[filtered_df['channel'] == channel]

        # Get data for previous and selected periods for the current channel
        before_data = channel_df[channel_df['month'] == previous_period.strftime('%Y-%m')].iloc[0, 1:].to_dict()
        last_data = channel_df[channel_df['month'] == selected_period.strftime('%Y-%m')].iloc[0, 1:].to_dict() \
            if not channel_df[channel_df['month'] == selected_period.strftime('%Y-%m')].empty else {}

        # Remove 'channel' from before_data and last_data
        if 'channel' in before_data:
            del before_data['channel']
        if 'channel' in last_data:
            del last_data['channel']

        # Create evaluation DataFrame for the current channel
        evaluation_media = pd.DataFrame(columns=['Channel', 'Metric', previous_period_name, selected_period_name])

        for metric, before_value in before_data.items():
            last_value = last_data.get(metric)  # Get the value for the last period, if available
            evaluation_media = pd.concat([evaluation_media, pd.DataFrame(
                {'Channel': [channel], 'Metric': [metric], previous_period_name: [before_value],
                 selected_period_name: [last_value]})], ignore_index=True)

        all_channel_evaluations.append(evaluation_media)

    # Concatenate all channel evaluation DataFrames
    final_evaluation = pd.concat(all_channel_evaluations, ignore_index=True)
    final_evaluation['Change %'] = round(
        ((final_evaluation[selected_period_name] - final_evaluation[previous_period_name]) /
         final_evaluation[previous_period_name]) * 100, 2
    )
    final_evaluation['Change %'] = pd.to_numeric(final_evaluation['Change %'], errors='coerce').fillna(0).round(
        2).astype(str) + '%'
    return final_evaluation


def get_logo_path_by_username(username):
    logo_map = {
        "alva": "template/logo/vinfast.svg.png",
        # "bardi": "assets/logo_bardi.png",
        # "sanken": "assets/logo_sanken.png"
    }
    return logo_map.get(username.lower(), "template/logo/default_logo.png")

def add_logo_to_title_slide(prs, logo_path, left=Inches(8), top=Inches(0.2), width=Inches(1)):
    title_slide = prs.slides[0]
    title_slide.shapes.add_picture(logo_path, left, top, width=width)

def add_footer_to_slide(slide):
    footer_text = f"Generated by AI Automation - Python | {date.today()}"

    # Ukuran dan posisi
    left = Inches(0.3)
    top = Inches(5.2)  # dekat bawah, tweak sesuai tinggi slide kamu
    width = Inches(9)
    height = Inches(0.3)

    textbox = slide.shapes.add_textbox(left, top, width, height)
    text_frame = textbox.text_frame
    text_frame.clear()

    p = text_frame.paragraphs[0]
    run = p.add_run()
    run.text = footer_text
    run.font.size = Pt(8)
    run.font.italic = True


def insert_table_from_placeholder(slide, df, placeholder_text):
    # Temukan shape yang punya placeholder
    target_shape = next((s for s in slide.shapes if s.has_text_frame and placeholder_text in s.text), None)
    if not target_shape:
      #  st.warning(f"Placeholder '{placeholder_text}' tidak ditemukan.")
        return

    # Ambil ukuran dan posisi
    left, top, width, height = target_shape.left, target_shape.top, target_shape.width, target_shape.height

    # Hapus shape placeholder
    target_shape._element.getparent().remove(target_shape._element)

    # Buat tabel
    rows, cols = df.shape
    table = slide.shapes.add_table(rows + 1, cols, left, top, width, height).table

    # Header
    for col_idx, col_name in enumerate(df.columns):
        cell = table.cell(0, col_idx)
        cell.text = str(col_name)
        cell.text_frame.paragraphs[0].font.bold = True

    # Data
    for row_idx, row_data in enumerate(df.values):
        for col_idx, value in enumerate(row_data):
            table.cell(row_idx + 1, col_idx).text = str(value)

    # Set font size
    for row in table.rows:
        for cell in row.cells:
            cell.text_frame.paragraphs[0].font.size = Pt(10)
def insert_content_from_placeholder(slide, df, placeholder_text, placeholder_type="text"):
    """
    Fungsi ini akan mengganti placeholder dengan tipe yang sesuai (teks, tabel, atau chart).
    - placeholder_type: "text", "table", atau "chart".
    """
    # Cari shape yang sesuai dengan placeholder
    placeholder_shape = next((s for s in slide.shapes if hasattr(s, "text") and placeholder_text in s.text), None)

    if placeholder_shape:
        if placeholder_type == "text":
            # Ganti placeholder dengan teks
            placeholder_shape.text = placeholder_text.replace("{text}", "Ganti dengan teks sesuai data")

        elif placeholder_type == "table":
            # Menambahkan tabel berdasarkan DataFrame
            rows, cols = df.shape
            left = Inches(2)
            top = Inches(2)

            # Menambahkan tabel
            table = slide.shapes.add_table(rows + 1, cols, left, top, Inches(6), Inches(0.8 + 0.4 * rows)).table

            # Menambahkan header tabel
            for col_idx, col_name in enumerate(df.columns):
                table.cell(0, col_idx).text = col_name

            # Menambahkan data ke tabel
            for row_idx, row in df.iterrows():
                for col_idx, value in enumerate(row):
                    table.cell(row_idx + 1, col_idx).text = str(value)

        elif placeholder_type == "chart":
            # Menambahkan chart (grafik batang)
            chart_data = CategoryChartData()
            chart_data.categories = df.columns.tolist()  # Menjadikan kolom sebagai kategori chart
            chart_data.add_series("Data", df.iloc[0].tolist())  # Ambil data pertama sebagai contoh series

            # Menambahkan chart di slide
            x, y, cx, cy = Inches(1), Inches(2), Inches(6), Inches(3)
            chart = slide.shapes.add_chart(
                XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, cx, cy, chart_data
            ).chart
def format_selected_period_id(period_id):
    # Contoh: "2025-01" → "Januari 2025"
    import calendar
    year, month = map(int, period_id.split('-'))
    bulan = calendar.month_name[month]
    return f"{bulan} {year}"
def insert_chart_from_placeholder(slide, img_stream, placeholder_text):
    """
    Gantikan placeholder_text dengan gambar chart dari BytesIO.
    """
    # Cari shape placeholder
    ph = next((s for s in slide.shapes
               if s.has_text_frame and placeholder_text in s.text), None)
    if not ph:
        return

    # Ambil posisi & ukuran
    left, top, width, height = ph.left, ph.top, ph.width, ph.height

    # Hapus placeholder
    ph._element.getparent().remove(ph._element)

    # Sisipkan gambar langsung dari stream
    slide.shapes.add_picture(img_stream, left, top, width=width, height=height)

def replace_text_placeholders(slide, content_dict):
    for shape in slide.shapes:
        if shape.has_text_frame:
            for placeholder, content in content_dict.items():
                if content["type"] == "text" and f"{{{placeholder}}}" in shape.text:
                    shape.text = shape.text.replace(f"{{{placeholder}}}", str(content["value"]))
def create_ppt_from_template(content_dict, template_path, username, selected_period):
    prs = Presentation(template_path)

    # Tambahkan logo di title slide
    logo_path = get_logo_path_by_username(username)
    add_logo_to_title_slide(prs, logo_path)

    # Slide 1: ubah brand dan periode
    slide_1 = prs.slides[0]
    for shape in slide_1.shapes:
        if hasattr(shape, "text"):
            shape.text = shape.text.replace("{brand}", username)
            shape.text = shape.text.replace("{periode}", format_selected_period_id(selected_period))

    # Proses semua slide
    for slide in prs.slides:
        add_footer_to_slide(slide)
        replace_text_placeholders(slide, content_dict)

        # Loop semua placeholder dari content_dict
        for placeholder, content in content_dict.items():
            if content["type"] == "table":
                if not content["value"].empty:
                    insert_table_from_placeholder(slide, content["value"], placeholder)
            elif content["type"] == "chart":
                # content["value"] sekarang adalah BytesIO
                insert_chart_from_placeholder(slide, content["value"], placeholder)
            elif content["type"] == "text":
                replace_text_placeholders(slide, {placeholder: content})  # Panggil fungsi untuk mengganti teks

    # Simpan ke memory
    ppt_bytes = BytesIO()
    prs.save(ppt_bytes)
    ppt_bytes.seek(0)
    return ppt_bytes
def generate_ppt(username, selected_period):
    # fact_post_ig = load_csv("data/fact_instagram_post_dummy.csv")
    fact_post_tt = load_csv("fact_tiktok_post.csv")
    fact_post_ig = load_csv("fact_instagram_post.csv")
    datamart_ig = load_csv("datamart_ig.csv")  #adjust disini
    datamart_tt = load_csv("datamart_tt.csv") #adjust disini
    datamart_all = pd.concat([datamart_ig, datamart_tt], ignore_index=True) #adjust disini
    datamart_all = datamart_all.groupby('month', as_index=False).sum()
    datamart_eval = load_csv("evaluation_media.csv")
    datamart_eval_all = pd.concat([datamart_ig, datamart_tt], ignore_index=True)
    #ig
    table_summary_ig = table_summary_channel(datamart_ig, selected_period)
    table_followers_trend_ig = table_followers_trend(datamart_ig, selected_period)
    table_engagement_trend_ig = table_engagement_trend(datamart_ig, selected_period)
    table_top_post_ig =table_top_post(fact_post_ig,selected_period)
    table_bottom_post_ig =table_bottom_post(fact_post_ig,selected_period)
    table_popular_post_ig = table_popular_post(fact_post_ig, selected_period)
    #tiktok
    table_summary_tt = table_summary_channel(datamart_tt,selected_period)
    table_followers_trend_tt = table_followers_trend(datamart_tt, selected_period)
    table_engagement_trend_tt = table_engagement_trend(datamart_tt, selected_period)
    table_top_post_tt =table_top_post(fact_post_tt,selected_period)
    table_bottom_post_tt =table_bottom_post(fact_post_tt,selected_period)
    table_popular_post_tt = table_popular_post(fact_post_tt, selected_period)
    #all
    table_summary_all = table_summary_channel(datamart_all,selected_period)
    #evaluation
    table_eval_media = table_evaluation_media(datamart_eval, selected_period)
    table_eval_socmed = table_evaluation_socmed(datamart_eval_all,selected_period)

    data_dict = {
        # "{table_summary}": table_summary_all,
        # "{table_summary_tiktok}": table_summary_tt,
        # "{table_followers_trend_tiktok}": table_followers_trend_tt,
        # "{table_engagement_trend_tiktok}": table_engagement_trend_tt,
        # "{table_popular_video_tiktok}":table_popular_post_tt,
        # "{table_summary_instagram}": table_summary_ig,
        # "{table_followers_trend_instagram}": table_followers_trend_ig,
        # "{table_engagement_trend_instagram}": table_engagement_trend_ig,
        # "{table_popular_post_instagram}": table_popular_post_ig,
        # "{text_top_post_instagram}" : table_top_post_ig,
        # "{table_bottom_post_instagram}" :table_bottom_post_ig

    }
    content_dict = {
        "text_top_post_instagram": {
            "type": "text",
            "value": generate_text_top_post(table_top_post_ig)
        },

        "text_bottom_post_instagram": {
            "type": "text",
            "value": generate_text_top_post(table_bottom_post_ig)
        },

        "text_popular_post_instagram": {
            "type": "text",
            "value": generate_text_top_post(table_popular_post_ig)
        },
        "text_top_post_tiktok": {
            "type": "text",
            "value": generate_text_top_post(table_top_post_tt)
        },

        "text_bottom_post_tiktok": {
            "type": "text",
            "value": generate_text_top_post(table_bottom_post_tt)
        },

        "text_popular_post_tiktok": {
            "type": "text",
            "value": generate_text_top_post(table_popular_post_tt)
        },
        "chart_engagement_trend_instagram": {
            "type": "chart",
            "value": plot_trend_engagement_io(table_engagement_trend_ig)
        },
        "chart_followers_trend_instagram": {
            "type": "chart",
            "value": plot_trend_followers_io(table_followers_trend_ig)
        },
        "chart_engagement_trend_tiktok": {
            "type": "chart",
            "value": plot_trend_engagement_io(table_engagement_trend_tt)
        },
        "chart_followers_trend_tiktok": {
            "type": "chart",
            "value": plot_trend_followers_io(table_followers_trend_tt)
        },
        "total_followers_all": {
            "type": "text",
            "value": table_summary_all['followers'].iloc[0]
        },
        "growth_followers_all": {
            "type": "text",
            "value": table_summary_all['followers_growth'].iloc[0]
        },
        "total_post_all": {
            "type": "text",
            "value": table_summary_all['total_post'].iloc[0]
        },
        "growth_post_all": {
            "type": "text",
            "value": table_summary_all['total_post_growth'].iloc[0]
        },
        "total_reach_all": {
            "type": "text",
            "value": table_summary_all['reach'].iloc[0]
        },
        "growth_reach_all": {
            "type": "text",
            "value": table_summary_all['reach_growth'].iloc[0]
        },
        "table_summary_tiktok": {
            "type": "table",
            "value": table_summary_tt
        },
        "table_summary_instagram": {
            "type": "table",
            "value": table_summary_ig
        },
        "evaluation_socmed": {
            "type": "table",
            "value": table_eval_socmed
        },
        "evaluation_media": {
            "type": "table",
            "value": table_eval_media
        }
    }

    ppt_bytes = create_ppt_from_template(content_dict, "template/template_ppt/[POC-ALVA_VinFast] template.pptx" , username, selected_period)
    return ppt_bytes


def plot_trend_engagement_io(df):
    # Pastikan df['month'] sudah datetime
    df = df.copy()
    df['month'] = pd.to_datetime(df['month'])

    fig, ax = plt.subplots(figsize=(6, 4))
    ax.plot(df['month'], df['engagement'], marker='o', linewidth=2)
    ax.set_xticks(df['month'])
    ax.set_xticklabels(df['month'].dt.strftime('%b %Y'), rotation=45)
    ax.set_title("Trend Engagement by Month")
    ax.set_xlabel("Month")
    ax.set_ylabel("Total Engagement")
    ax.grid(True)
    plt.tight_layout()

    # Dump ke BytesIO
    img_stream = io.BytesIO()
    fig.savefig(img_stream, format="png", bbox_inches="tight")
    plt.close(fig)
    img_stream.seek(0)
    return img_stream

def plot_trend_followers_io(df):
    # Pastikan df['month'] sudah datetime
    df = df.copy()
    df['month'] = pd.to_datetime(df['month'])

    fig, ax = plt.subplots(figsize=(6, 4))
    ax.plot(df['month'], df['followers'], marker='o', linewidth=2)
    ax.set_xticks(df['month'])
    ax.set_xticklabels(df['month'].dt.strftime('%b %Y'), rotation=45)
    ax.set_title("Trend Followers by Month")
    ax.set_xlabel("Month")
    ax.set_ylabel("Total Followers")
    ax.grid(True)
    plt.tight_layout()

    # Dump ke BytesIO
    img_stream = io.BytesIO()
    fig.savefig(img_stream, format="png", bbox_inches="tight")
    plt.close(fig)
    img_stream.seek(0)
    return img_stream
