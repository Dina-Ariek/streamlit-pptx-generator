import streamlit as st
import pandas as pd
from pptx import Presentation
from io import BytesIO
from pptx.util import Pt
def load_instagram_post(file_path):
    """
    Load data fact_instagram_post dari file CSV.
    :param file_path: path ke file fact_instagram_post.csv
    :return: DataFrame
    """
    df = pd.read_csv(file_path)
    return df


def load_tiktok_post(file_path):
    """
    Load data fact_tiktok_post dari file CSV.
    :param file_path: path ke file fact_tiktok_post.csv
    :return: DataFrame
    """
    df = pd.read_csv(file_path)
    return df


def aggregate_instagram_monthly(profile_df: pd.DataFrame, post_df: pd.DataFrame) -> pd.DataFrame:
    # Pastikan format datetime
    profile_df['date'] = pd.to_datetime(profile_df['date'])
    post_df['post_date'] = pd.to_datetime(post_df['post_date'])

    # Buat kolom Month
    profile_df['Month'] = profile_df['date'].dt.to_period('M').astype(str)
    post_df['Month'] = post_df['post_date'].dt.to_period('M').astype(str)

    # Hitung Growth Followers per bulan
    growth_df = profile_df.groupby('Month').agg({
        'followers': ['first', 'last']
    })
    growth_df.columns = ['followers_start', 'followers_end']
    growth_df['Growth'] = growth_df['followers_end'] - growth_df['followers_start']

    # Agregasi Profile Daily
    profile_agg = profile_df.groupby('Month').agg({
        'profile_reach': 'sum',
        'profile_visit': 'sum',
        'followers': 'mean'
    }).reset_index()

    # Agregasi Post Performance
    post_agg = post_df.groupby('Month').agg({
        'like_count': 'sum',
        'comment_count': 'sum',
        'share_count': 'sum',
        'saved_count': 'sum'
    }).reset_index()

    #engagement
    post_agg['engagement'] = post_agg['like_count'] + post_agg['comment_count'] + post_agg['share_count'] + post_agg['saved_count']

    # Merge semua
    final_df = profile_agg.merge(growth_df['Growth'], on='Month') \
        .merge(post_agg, on='Month')

    # Hitung ER Reach (%) dan ER Followers (%)
    final_df['ER Reach (%)'] = round((final_df['engagement'] / final_df['profile_reach']) * 100, 2)
    final_df['ER Followers (%)'] = round((final_df['engagement'] / final_df['followers']) * 100, 2)

    return final_df

def table_performance_ig(ig_monthly_df, selected_period):
    df = pd.DataFrame(ig_monthly_df)
    # Hitung Gap (hanya untuk kolom numerik tertentu)
    gap_row = (df.iloc[1, 1:] - df.iloc[0, 1:]) / df.iloc[0, 1:]

    # Handle kolom ER Reach dan ER Followers jadi NaN
    gap_row[['ER Reach (%)', 'ER Followers (%)']] = ''

    # Tambahkan kolom Month = 'Gap'
    gap_row['Month'] = 'Gap'

    # Urutkan kolom sesuai df awal
    gap_row = gap_row[df.columns]

    # Gabungkan ke dataframe awal
    df = pd.concat([df, pd.DataFrame([gap_row])], ignore_index=True)

    return(df)

def create_ppt_from_template(data, template_path, username):
    """Membuat file PPT berdasarkan template dan data yang diunggah."""
    prs = Presentation(template_path)
    header_text = f"Prepared by {username}  |   Social Media Report @2022-2024"
    for slide in prs.slides:
        # Cek apakah header sudah ada di template (gunakan placeholder)
        header_shape = next((s for s in slide.shapes if s.has_text_frame and "{header}" in s.text), None)

        if header_shape:
            header_shape.text = header_text  # Ganti teks placeholder


    #Slide 1
    slide_1 = prs.slides[0]
    for shape in slide_1.shapes:
        if hasattr(shape, "text"):
            shape.text = shape.text.replace("{product_name}", username)
            shape.text = shape.text.replace("{periode}", "March 2024")

    slide_2 = prs.slides[1]
    for shape in slide_2.shapes:
        if hasattr(shape, "text"):
            shape.text = shape.text.replace("{subtitle}", "Instagram")
    slide = prs.slides[2]

    # Temukan textbox placeholder
    target_shape = next((s for s in slide.shapes if s.has_text_frame and "{tabel}" in s.text), None)
    if not target_shape:
        st.error("⚠ Placeholder tabel tidak ditemukan!")
        return None

    # Ambil ukuran textbox, lalu hapus
    left, top, width, height = target_shape.left, target_shape.top, target_shape.width, target_shape.height
    target_shape._element.getparent().remove(target_shape._element)

    # Buat tabel
    rows, cols = data.shape
    table = slide.shapes.add_table(rows + 1, cols, left, top, width, height).table

    # Set lebar kolom dan tinggi baris
    for col in range(cols): table.columns[col].width = width // cols
    for row in range(rows + 1): table.rows[row].height = height // (rows + 1)

    # Isi header
    for col_idx, col_name in enumerate(data.columns):
        cell = table.cell(0, col_idx)
        cell.text, cell.text_frame.paragraphs[0].font.bold = col_name, True

    # Isi data & atur font size ke 10 pt
    for row_idx, row_data in enumerate(data.values):
        for col_idx, value in enumerate(row_data):
            table.cell(row_idx + 1, col_idx).text = str(value)

    # Atur semua font jadi 10 pt
    for row in table.rows:
        for cell in row.cells:
            cell.text_frame.paragraphs[0].font.size = Pt(10)

    # Simpan PPT ke BytesIO
    ppt_bytes = BytesIO()
    prs.save(ppt_bytes)
    ppt_bytes.seek(0)
    return ppt_bytes

def generate_ppt(username, selected_period):
    template_path = "template-alva.pptx"
    username = username
    ig_post_df = load_instagram_post(r'data/fact_instagram_post_dummy.csv')
    ig_profile_df = pd.read_csv(r'data/fact_instagram_profile_daily_dummy.csv')
    ig_monthly_df = aggregate_instagram_monthly(ig_profile_df, ig_post_df)
    table_perform_ig = table_performance_ig(ig_monthly_df,selected_period)
    pptx_file = create_ppt_from_template(table_perform_ig, template_path, username)
    return pptx_file
# ig_post_df = load_instagram_post(r'data/fact_instagram_post_dummy.csv')

# with open("output.pptx", "wb") as f:
#     f.write(pptx_file.getbuffer())



